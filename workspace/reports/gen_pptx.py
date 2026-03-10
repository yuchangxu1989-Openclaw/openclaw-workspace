#!/usr/bin/env python3
"""Generate OpenClaw Product Insight V3b PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x0D, 0x23, 0x3E)
MID_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
LIGHT_BLUE = RGBColor(0x2D, 0x5F, 0x8A)
ORANGE = RGBColor(0xF0, 0x82, 0x1D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        sld = shape.fill._fill
        solidFill = sld.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                a = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                a.set('val', str(int(alpha * 1000)))
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, line_space=1.2):
    """Add textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.3)
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=18, color=DARK_GRAY, bullet_color=ORANGE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(font_size * 0.4)
    return txBox

def add_big_number(slide, left, top, width, height, number, label, num_color=ORANGE):
    add_text(slide, left, top, width, Inches(1), number, font_size=54, color=num_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(1.1), width, Inches(0.6), label, font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ========== SLIDE 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(0), Inches(0), W, H, DARK_BLUE)
# Accent bar
add_rect(slide, Inches(1), Inches(2.8), Inches(1.5), Pt(4), ORANGE)
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5), 'OpenClaw 产品洞察', font_size=48, color=WHITE, bold=True)
add_text(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.8), '从"AI工具"到"AI同事"的范式跃迁', font_size=28, color=RGBColor(0xCC, 0xCC, 0xCC))
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), '版本 V3b ｜ 2026年3月 ｜ 高管汇报级', font_size=16, color=RGBColor(0x88, 0x99, 0xAA))
add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5), 'Claude Opus 4 · coder', font_size=14, color=RGBColor(0x66, 0x77, 0x88))

# ========== SLIDE 2: TOC ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(1), Inches(0.8), Inches(10), Inches(0.8), '目  录', font_size=36, color=DARK_BLUE, bold=True)
add_rect(slide, Inches(1), Inches(1.6), Inches(2), Pt(3), ORANGE)

toc_items = [
    ('01', '产品定位', '从"AI工具"到"AI同事"的范式跃迁'),
    ('02', '场景分析', '从开发者玩具到生活基础设施'),
    ('03', '用户粘性与社区', '为什么用过的人回不去了'),
    ('04', '能力边界', 'LLM的"物理学"与"外骨骼"策略'),
    ('05', '未来展望', '路线图与关键风险'),
    ('06', '对vivo的启示', '"AI硬件伴侣"——被忽视的路径'),
]
for i, (num, title, sub) in enumerate(toc_items):
    y = Inches(2.2 + i * 0.8)
    add_text(slide, Inches(1.2), y, Inches(0.8), Inches(0.6), num, font_size=28, color=ORANGE, bold=True)
    add_text(slide, Inches(2.2), y, Inches(3), Inches(0.4), title, font_size=22, color=DARK_BLUE, bold=True)
    add_text(slide, Inches(2.2), y + Inches(0.35), Inches(8), Inches(0.4), sub, font_size=14, color=MED_GRAY)

# ========== SECTION 1: Product Positioning (Slides 3-6) ==========

# Slide 3: Section cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(1), Inches(2), Inches(0.6), '01', font_size=60, color=ORANGE, bold=True)
add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2), '产品定位\n从"AI工具"到"AI同事"的范式跃迁', font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(4.0), Inches(1.5), Pt(4), ORANGE)
add_text(slide, Inches(1), Inches(4.3), Inches(11), Inches(1.5),
    'OpenClaw不是又一个AI聊天机器人，而是全球首个\n"个人AI操作系统网关"——让AI从被动工具变成主动同事',
    font_size=20, color=RGBColor(0xBB, 0xCC, 0xDD))

# Slide 4: Four-quadrant positioning
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '精确定位：四象限坐标系', font_size=30, color=DARK_BLUE, bold=True)

# Quadrant boxes
quads = [
    (Inches(0.8), Inches(1.5), '开发者工具', 'Claude Code / Cursor / Copilot\n需要编程能力\n服务代码生产场景', MID_BLUE),
    (Inches(7), Inches(1.5), '普通用户工具', 'ChatGPT / Gemini\n对话界面\n功能限于问答，无法"做事"', MID_BLUE),
    (Inches(0.8), Inches(4.0), '编排框架/自主实验', 'LangChain / CrewAI / AutoGPT\n面向开发者的构建工具\n目标宏大但可靠性不足', MID_BLUE),
    (Inches(7), Inches(4.0), '🟠 OpenClaw', '既有聊天渠道接入 · 零学习成本\n24/7持久记忆 · 主动行为\n邮件+日历+设备+文件+浏览器', ORANGE),
]
for x, y, title, desc, color in quads:
    box = add_rect(slide, x, y, Inches(5.3), Inches(2.2), LIGHT_GRAY if color != ORANGE else RGBColor(0xFF, 0xF3, 0xE0))
    box.line.color.rgb = color
    box.line.width = Pt(2)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(4.7), Inches(0.5), title, font_size=20, color=color, bold=True)
    add_multi_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(4.7), Inches(1.3), desc.split('\n'), font_size=15, color=DARK_GRAY)

# Slide 5: Three "Firsts"
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '三个"第一次"——变革性论证', font_size=30, color=DARK_BLUE, bold=True)

firsts = [
    ('自然存在感', '通过WhatsApp/Telegram/Discord等既有聊天渠道\n消除"打开AI工具"这个动作\nAI成为通讯录里的"联系人"'),
    ('连续自我', 'MEMORY.md + SOUL.md 持久记忆与人格\n用户关系从工具调用变为日积月累的"相处"\n切换成本 ≈ 换掉一个懂你的人'),
    ('行动力', '心跳机制 · cron定时任务 · 后台Agent\n不在时主动：检查邮件、监控日程、响应事件\n从"响应式助手"到"主动式同事"'),
]
for i, (title, desc) in enumerate(firsts):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(1.5), Inches(3.7), Inches(5.0), LIGHT_GRAY)
    add_rect(slide, x, Inches(1.5), Inches(3.7), Pt(4), ORANGE)
    add_text(slide, x + Inches(0.3), Inches(1.8), Inches(3.1), Inches(0.6), f'第一次', font_size=14, color=ORANGE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(2.2), Inches(3.1), Inches(0.6), title, font_size=24, color=DARK_BLUE, bold=True)
    add_multi_text(slide, x + Inches(0.3), Inches(3.0), Inches(3.1), Inches(3.0), desc.split('\n'), font_size=15, color=DARK_GRAY)

# Slide 6: Community voices
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '社区真实声音', font_size=30, color=DARK_BLUE, bold=True)

quotes = [
    ('"Setup OpenClaw yesterday. All I have to say is, wow...\nThe future is already here."', '@jonahships_'),
    ('"Even if LLMs suddenly stopped improving, we could spend\nyears discovering new transformative uses. OpenClaw feels\nlike that kind of leap forward."', '@markjaquith'),
    ('"Personal AI assistant undersells it — it\'s a company assistant,\nfamily assistant, team tool. Proactive AF. Memory is amazing,\ncontext persists 24/7."', '@danpeguine'),
]
for i, (quote, author) in enumerate(quotes):
    y = Inches(1.3 + i * 1.9)
    add_rect(slide, Inches(1.2), y, Pt(4), Inches(1.5), ORANGE)
    add_multi_text(slide, Inches(1.6), y + Inches(0.1), Inches(9), Inches(1.2), quote.split('\n'), font_size=17, color=DARK_GRAY)
    add_text(slide, Inches(1.6), y + Inches(1.3), Inches(5), Inches(0.3), f'—— {author}，X/Twitter 2026', font_size=13, color=MED_GRAY)

# ========== SECTION 2: Scenarios (Slides 7-9) ==========

# Slide 7: Section cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(1), Inches(2), Inches(0.6), '02', font_size=60, color=ORANGE, bold=True)
add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2), '场景分析\n从开发者玩具到生活基础设施', font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(4.0), Inches(1.5), Pt(4), ORANGE)
add_text(slide, Inches(1), Inches(4.3), Inches(11), Inches(1),
    '关键转折点：当安装一个技能像安装App一样简单时\n场景将指数级扩展', font_size=20, color=RGBColor(0xBB, 0xCC, 0xDD))

# Slide 8: Current scenario distribution
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '当前使用场景分布（推测）', font_size=30, color=DARK_BLUE, bold=True)

scenarios = [
    ('开发/编程辅助', 35, 'Claude Code会话管理、代码审查、自动化测试'),
    ('个人信息管理', 25, '邮件管理、日历调度、航班值机、待办追踪'),
    ('知识管理与研究', 15, '网页抓取、信息整理、文档摘要'),
    ('智能家居/设备', 10, '通过Node机制控制Mac/树莓派/智能设备'),
    ('创意与内容', 10, '写作辅助、TTS语音生成、图片处理'),
    ('团队协作', 5, 'Discord/Slack群组中的团队AI助手'),
]
max_bar_w = Inches(6)
for i, (name, pct, desc) in enumerate(scenarios):
    y = Inches(1.4 + i * 0.9)
    add_text(slide, Inches(0.8), y, Inches(2.5), Inches(0.4), name, font_size=16, color=DARK_BLUE, bold=True)
    bar_w = int(max_bar_w * pct / 35)
    add_rect(slide, Inches(3.5), y + Inches(0.05), bar_w, Inches(0.35), ORANGE if pct >= 25 else LIGHT_BLUE)
    add_text(slide, Inches(3.5) + bar_w + Inches(0.2), y, Inches(1), Inches(0.4), f'{pct}%', font_size=16, color=ORANGE if pct >= 25 else LIGHT_BLUE, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.4), Inches(8), Inches(0.3), desc, font_size=12, color=MED_GRAY)

add_text(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.4),
    '注：基于公开社交媒体帖子(50-80条)+产品功能页的定性推测，存在幸存者偏差', font_size=11, color=MED_GRAY)

# Slide 9: Five trends
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '五大长期趋势', font_size=30, color=DARK_BLUE, bold=True)

trends = [
    ('非技术化迁移', '从"帮我写代码"到"替我管生活"'),
    ('社会单元扩展', '从"单人助手"到"家庭/团队AI"'),
    ('交互消隐', '从"聊天交互"到"环境智能"'),
    ('Multi-Agent协同', '从一个AI到一支AI团队'),
    ('技能生态', '从封闭能力到开放市场(App Store模式)'),
]
for i, (title, desc) in enumerate(trends):
    x = Inches(0.6 + (i % 5) * 2.45)
    y = Inches(1.5)
    h = Inches(4.5)
    add_rect(slide, x, y, Inches(2.2), h, LIGHT_GRAY)
    add_rect(slide, x, y, Inches(2.2), Pt(4), ORANGE)
    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(1.8), Inches(0.5), f'趋势{i+1}', font_size=13, color=ORANGE, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.7), Inches(1.8), Inches(0.8), title, font_size=20, color=DARK_BLUE, bold=True)
    add_multi_text(slide, x + Inches(0.2), y + Inches(1.5), Inches(1.8), Inches(2.5), [desc], font_size=14, color=DARK_GRAY)

# Slide 10: Scenario diffusion path
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '场景扩散路径', font_size=30, color=DARK_BLUE, bold=True)

phases = [
    ('当前-2026H2', '阶段1', '技术早期采用者\n编程 + 自动化', LIGHT_BLUE),
    ('2026H2-2027', '阶段2', '高价值职业用户\n信息管理 + 内容', MID_BLUE),
    ('2027-2028', '阶段3', '家庭与团队\n多用户 + 设备联动', DARK_BLUE),
    ('2028+', '阶段4', '大众基础设施\n技能市场成熟', ORANGE),
]
for i, (time, phase, desc, color) in enumerate(phases):
    x = Inches(0.8 + i * 3.1)
    # Arrow-like shape
    add_rect(slide, x, Inches(2.5), Inches(2.8), Inches(3.5), color)
    add_text(slide, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.4), time, font_size=14, color=WHITE)
    add_text(slide, x + Inches(0.2), Inches(3.2), Inches(2.4), Inches(0.5), phase, font_size=24, color=WHITE, bold=True)
    add_multi_text(slide, x + Inches(0.2), Inches(4.0), Inches(2.4), Inches(1.5), desc.split('\n'), font_size=16, color=WHITE)

# ========== SECTION 3: User Stickiness (Slides 11-13) ==========

# Slide 11: Section cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(1), Inches(2), Inches(0.6), '03', font_size=60, color=ORANGE, bold=True)
add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2), '用户粘性与社区\n为什么用过的人回不去了', font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(4.0), Inches(1.5), Pt(4), ORANGE)
add_text(slide, Inches(1), Inches(4.3), Inches(11), Inches(1),
    '粘性不来自功能强大，而来自"关系深度"\n切换成本 ≈ 换掉一个了解你的同事', font_size=20, color=RGBColor(0xBB, 0xCC, 0xDD))

# Slide 12: Five core traits
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '5个核心粘性特质', font_size=30, color=DARK_BLUE, bold=True)

traits = [
    ('🏠', '自然存在感', '在聊天App里\n= 一直在那里\n= 默认依赖'),
    ('🧠', '记忆连续性', '信任 =\n被理解 × 时间积累'),
    ('⚡', '主动执行力', '主动性是区分\n"工具"和"同事"\n的分水岭'),
    ('🔒', '数据主权', '控制感\n是信任的前提\n自托管 = 安全感'),
    ('🎭', '个性化深度', '个性化 =\n不可替代性'),
]
for i, (icon, title, desc) in enumerate(traits):
    x = Inches(0.5 + i * 2.5)
    add_rect(slide, x, Inches(1.5), Inches(2.2), Inches(4.5), LIGHT_GRAY)
    add_text(slide, x, Inches(1.7), Inches(2.2), Inches(0.8), icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.5), Inches(1.8), Inches(0.6), title, font_size=20, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.2), Inches(3.3), Inches(1.8), Inches(2.5), desc.split('\n'), font_size=15, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Slide 13: Switching cost
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '切换成本构成', font_size=30, color=DARK_BLUE, bold=True)

costs = [
    ('记忆迁移', '日积月累的上下文无法导出', '极高'),
    ('技能迁移', '定制技能 = 专属能力库', '高'),
    ('习惯重建', '交互习惯+工作流程重构', '中高'),
    ('社交迁移', '群组中的AI身份/人格', '中'),
    ('情感沉没', '类似"换掉一个老朋友"', '极高'),
]
for i, (name, desc, level) in enumerate(costs):
    y = Inches(1.3 + i * 1.1)
    add_rect(slide, Inches(1), y, Inches(11), Inches(0.9), LIGHT_GRAY)
    add_text(slide, Inches(1.3), y + Inches(0.1), Inches(2), Inches(0.4), name, font_size=20, color=DARK_BLUE, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.15), Inches(5), Inches(0.4), desc, font_size=16, color=DARK_GRAY)
    level_color = RGBColor(0xCC, 0x33, 0x33) if '极高' in level else ORANGE
    add_text(slide, Inches(9.5), y + Inches(0.15), Inches(2), Inches(0.4), f'🔴 {level}', font_size=18, color=level_color, bold=True)

# ========== SECTION 4: Capability Boundaries (Slides 14-16) ==========

# Slide 14: Section cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(1), Inches(2), Inches(0.6), '04', font_size=60, color=ORANGE, bold=True)
add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2), '能力边界\nLLM的"物理学"与"外骨骼"策略', font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(4.0), Inches(1.5), Pt(4), ORANGE)
add_text(slide, Inches(1), Inches(4.3), Inches(11), Inches(1),
    '赢家不是"模型最强的"，而是"外骨骼最好的"', font_size=22, color=RGBColor(0xBB, 0xCC, 0xDD))

# Slide 15: LLM capability metrics
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), 'LLM能力边界量化（Claude Opus 4, 200+用例）', font_size=26, color=DARK_BLUE, bold=True)

metrics = [
    ('执行链步骤', '≤4步 @95%', '7+步 <40%'),
    ('跨Turn记忆', '同Turn 90%', '2+次 <10%'),
    ('并行追踪', '≤3 可靠', '7+ 严重丢失'),
    ('规则遵从', '≤5条 @80%', '161条 <5%'),
    ('自检能力', '表面 80%', '系统性 10%'),
    ('意图推理', 'IC1-2 @95%', 'IC5 @30%'),
]
for i, (name, good, bad) in enumerate(metrics):
    y = Inches(1.3 + i * 0.95)
    add_text(slide, Inches(0.8), y + Inches(0.1), Inches(2.5), Inches(0.4), name, font_size=18, color=DARK_BLUE, bold=True)
    # Good
    add_rect(slide, Inches(3.5), y, Inches(3.5), Inches(0.7), RGBColor(0xE8, 0xF5, 0xE9))
    add_text(slide, Inches(3.7), y + Inches(0.15), Inches(3.1), Inches(0.4), f'✅ {good}', font_size=16, color=RGBColor(0x2E, 0x7D, 0x32))
    # Bad
    add_rect(slide, Inches(7.3), y, Inches(3.5), Inches(0.7), RGBColor(0xFF, 0xEB, 0xEE))
    add_text(slide, Inches(7.5), y + Inches(0.15), Inches(3.1), Inches(0.4), f'⚠️ {bad}', font_size=16, color=RGBColor(0xC6, 0x28, 0x28))

add_text(slide, Inches(0.8), Inches(7.0), Inches(12), Inches(0.3),
    '核心不等式：用户期望的任务复杂度 >> LLM裸奔的可靠区间', font_size=16, color=ORANGE, bold=True)

# Slide 16: Exoskeleton architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '"外骨骼"分层架构', font_size=30, color=DARK_BLUE, bold=True)

layers = [
    ('Layer 4', '意图理解', 'LLM驱动，允许模糊', ORANGE),
    ('Layer 3', '任务规划', 'LLM辅助 + 模板库，半确定性', RGBColor(0xE8, 0x8D, 0x2A)),
    ('Layer 2', '步骤执行', '程序化为主', MID_BLUE),
    ('Layer 1', '验证与回滚', '纯程序化，零LLM依赖', LIGHT_BLUE),
    ('Layer 0', '持久化', '纯确定性存储', DARK_BLUE),
]
for i, (layer, name, desc, color) in enumerate(layers):
    y = Inches(1.3 + i * 1.15)
    w = Inches(11 - i * 0.8)
    x = Inches(0.8 + i * 0.4)
    add_rect(slide, x, y, w, Inches(0.95), color)
    add_text(slide, x + Inches(0.3), y + Inches(0.1), Inches(1.5), Inches(0.4), layer, font_size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(2), y + Inches(0.1), Inches(2.5), Inches(0.4), name, font_size=18, color=WHITE, bold=True)
    add_text(slide, x + Inches(5), y + Inches(0.1), Inches(4), Inches(0.4), desc, font_size=15, color=WHITE)

add_text(slide, Inches(1.5), Inches(7.0), Inches(10), Inches(0.3),
    '↑ LLM依赖度高（允许模糊）          ↓ 确定性高（零容错）', font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ========== SECTION 5: Future (Slides 17-19) ==========

# Slide 17: Section cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(1), Inches(2), Inches(0.6), '05', font_size=60, color=ORANGE, bold=True)
add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2), '未来展望\n路线图与关键风险', font_size=36, color=WHITE, bold=True)

# Slide 18: Roadmap
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '发展路线图', font_size=30, color=DARK_BLUE, bold=True)

roadmap = [
    ('短期 2026H1-H2', '3-5个垂直场景极致可靠性\n技能库100+\n核心用户群建立', ORANGE),
    ('中期 2027', '技能市场上线\n家庭AI场景\n企业轻量版\n技能500+', MID_BLUE),
    ('长期 2028+', '个人AI操作系统\n事实标准\n生态成熟', DARK_BLUE),
]
for i, (period, items, color) in enumerate(roadmap):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(1.5), Inches(3.7), Inches(5.0), color)
    add_text(slide, x + Inches(0.3), Inches(1.8), Inches(3.1), Inches(0.5), period, font_size=22, color=WHITE, bold=True)
    add_rect(slide, x + Inches(0.3), Inches(2.5), Inches(1), Pt(3), WHITE)
    add_multi_text(slide, x + Inches(0.3), Inches(2.9), Inches(3.1), Inches(3.0), items.split('\n'), font_size=17, color=WHITE)

# Slide 19: Key risks
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '关键风险', font_size=30, color=DARK_BLUE, bold=True)

risks = [
    ('模型能力跳变', '中 30-40%', '高', '外骨骼转型为可靠性增强层'),
    ('大厂封闭生态', '高 60-70%', '高', '开放生态差异化'),
    ('技术门槛限制增长', '高 70%', '高', '一键安装 + 预配置硬件'),
]
# Header
y_h = Inches(1.3)
add_rect(slide, Inches(0.8), y_h, Inches(11.5), Inches(0.6), DARK_BLUE)
for col, (text, w) in enumerate([('风险', 3), ('概率', 2), ('影响', 1.5), ('缓解策略', 5)]):
    x_offset = [0.8, 3.8, 5.8, 7.3][col]
    add_text(slide, Inches(x_offset), y_h + Inches(0.1), Inches(w), Inches(0.4), text, font_size=16, color=WHITE, bold=True)

for i, (risk, prob, impact, mitigation) in enumerate(risks):
    y = Inches(2.1 + i * 1.2)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.0), bg_color)
    add_text(slide, Inches(1.0), y + Inches(0.25), Inches(2.8), Inches(0.5), risk, font_size=17, color=DARK_BLUE, bold=True)
    add_text(slide, Inches(3.8), y + Inches(0.25), Inches(2), Inches(0.5), prob, font_size=16, color=ORANGE, bold=True)
    add_text(slide, Inches(5.8), y + Inches(0.25), Inches(1.5), Inches(0.5), f'🔴 {impact}', font_size=16, color=RGBColor(0xCC, 0x33, 0x33))
    add_text(slide, Inches(7.3), y + Inches(0.25), Inches(5), Inches(0.5), mitigation, font_size=16, color=DARK_GRAY)

# ========== SECTION 6: vivo Insights (Slides 20-24) ==========

# Slide 20: Section cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(1), Inches(2), Inches(0.6), '06', font_size=60, color=ORANGE, bold=True)
add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2), '对vivo的启示\n"AI硬件伴侣"——被忽视的路径', font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(1), Inches(4.0), Inches(1.5), Pt(4), ORANGE)
add_text(slide, Inches(1), Inches(4.3), Inches(11), Inches(1),
    '不是在手机里加AI功能\n而是做独立的"AI硬件伴侣"', font_size=22, color=RGBColor(0xBB, 0xCC, 0xDD))

# Slide 21: Market data
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '市场规模——为什么现在必须行动', font_size=30, color=DARK_BLUE, bold=True)

add_big_number(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(2),
    '$182.97B', 'AI Agent市场 2033年规模')
add_big_number(slide, Inches(4.7), Inches(1.8), Inches(3.8), Inches(2),
    '49.6%', 'CAGR 年复合增长率')
add_big_number(slide, Inches(8.9), Inches(1.8), Inches(3.8), Inches(2),
    '$143.06B', '边缘AI市场 2034年规模')

add_rect(slide, Inches(1), Inches(4.2), Inches(11), Pt(1), MED_GRAY)

add_text(slide, Inches(1), Inches(4.5), Inches(5), Inches(0.4), 'AI Agent市场', font_size=18, color=DARK_BLUE, bold=True)
add_text(slide, Inches(1), Inches(5.0), Inches(5), Inches(0.3), '$7.63B (2025) → $182.97B (2033)', font_size=16, color=DARK_GRAY)
add_text(slide, Inches(1), Inches(5.4), Inches(5), Inches(0.3), 'CAGR 49.6%', font_size=16, color=ORANGE, bold=True)

add_text(slide, Inches(7), Inches(4.5), Inches(5), Inches(0.4), '边缘AI市场', font_size=18, color=DARK_BLUE, bold=True)
add_text(slide, Inches(7), Inches(5.0), Inches(5), Inches(0.3), '$25.65B (2025) → $143.06B (2034)', font_size=16, color=DARK_GRAY)
add_text(slide, Inches(7), Inches(5.4), Inches(5), Inches(0.3), 'CAGR 21.04%', font_size=16, color=ORANGE, bold=True)

add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.8),
    '⏰ 窗口期：12-18个月（推测）  |  💰 MVP投资：¥3000-5000万（推测）',
    font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Slide 22: vivo advantages
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), 'vivo独特优势', font_size=30, color=DARK_BLUE, bold=True)

advantages = [
    ('🔧', 'OriginOS\n系统级权限', '不可复制的\n底层能力'),
    ('🏭', '硬件设计\n供应链', '纯软件公司\n无法匹敌'),
    ('👥', '4亿+\n存量用户', '线下渠道\n触达能力'),
    ('🎯', '品牌\nAI空白', '定义窗口\n先发优势'),
]
for i, (icon, title, desc) in enumerate(advantages):
    x = Inches(0.8 + i * 3.1)
    add_rect(slide, x, Inches(1.5), Inches(2.8), Inches(4.5), LIGHT_GRAY)
    add_rect(slide, x, Inches(1.5), Inches(2.8), Pt(4), ORANGE)
    add_text(slide, x, Inches(1.8), Inches(2.8), Inches(0.8), icon, font_size=44, alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.3), Inches(2.8), Inches(2.2), Inches(1.2), title.split('\n'), font_size=22, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_multi_text(slide, x + Inches(0.3), Inches(4.2), Inches(2.2), Inches(1), desc.split('\n'), font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Slide 23: Two product concepts
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '两个产品概念', font_size=30, color=DARK_BLUE, bold=True)

# Product 1
add_rect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), RGBColor(0xFF, 0xF3, 0xE0))
add_rect(slide, Inches(0.8), Inches(1.3), Inches(5.5), Pt(5), ORANGE)
add_text(slide, Inches(1.2), Inches(1.6), Inches(4.8), Inches(0.5), '💼 AI办公搭子', font_size=26, color=ORANGE, bold=True)
add_text(slide, Inches(1.2), Inches(2.3), Inches(4.8), Inches(0.4), '桌面AI助手盒子', font_size=20, color=DARK_BLUE, bold=True)
add_text(slide, Inches(1.2), Inches(2.8), Inches(4.8), Inches(0.5), '¥999 - 2,999', font_size=28, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(3.5), Inches(4.8), Inches(3),
    ['日程管理与智能提醒', '邮件处理与自动回复', '会议纪要与任务分发', '文档搜索与知识管理'], font_size=16)

# Product 2
add_rect(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(5.5), RGBColor(0xE3, 0xF2, 0xFD))
add_rect(slide, Inches(7), Inches(1.3), Inches(5.5), Pt(5), DARK_BLUE)
add_text(slide, Inches(7.4), Inches(1.6), Inches(4.8), Inches(0.5), '🏠 AI家庭搭子', font_size=26, color=DARK_BLUE, bold=True)
add_text(slide, Inches(7.4), Inches(2.3), Inches(4.8), Inches(0.4), '家庭AI中枢', font_size=20, color=DARK_BLUE, bold=True)
add_text(slide, Inches(7.4), Inches(2.8), Inches(4.8), Inches(0.5), '¥1,499 - 3,999', font_size=28, color=DARK_BLUE, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(3.5), Inches(4.8), Inches(3),
    ['家庭日程协调与提醒', '智能家居设备联动', '孩子学习辅助与监护', '家庭信息枢纽'], font_size=16)

# Slide 24: Go/No-Go
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, DARK_BLUE)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.6), '三步走 + Go/No-Go 门禁', font_size=30, color=DARK_BLUE, bold=True)

steps = [
    ('P0', '0-6月', 'POC原型', ['NPS ≥ 50', '任务成功率 ≥ 85%'], ORANGE),
    ('P1', '6-12月', 'MVP灰度', ['30日留存 ≥ 40%', '退货率 < 15%'], MID_BLUE),
    ('P2', '12-24月', '规模化', ['年出货 ≥ 50万台', '生态技能 200+'], DARK_BLUE),
]
for i, (phase, time, name, gates, color) in enumerate(steps):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(1.3), Inches(3.7), Inches(5.5), color)
    add_text(slide, x + Inches(0.3), Inches(1.5), Inches(3.1), Inches(0.6), phase, font_size=36, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(2.2), Inches(3.1), Inches(0.4), time, font_size=18, color=WHITE)
    add_text(slide, x + Inches(0.3), Inches(2.7), Inches(3.1), Inches(0.5), name, font_size=24, color=WHITE, bold=True)
    add_rect(slide, x + Inches(0.3), Inches(3.4), Inches(1), Pt(2), WHITE)
    add_text(slide, x + Inches(0.3), Inches(3.7), Inches(3.1), Inches(0.4), 'Go/No-Go 门禁:', font_size=14, color=RGBColor(0xDD, 0xDD, 0xDD))
    for j, gate in enumerate(gates):
        add_text(slide, x + Inches(0.3), Inches(4.2 + j * 0.5), Inches(3.1), Inches(0.4),
            f'✓ {gate}', font_size=18, color=WHITE, bold=True)

# ========== SLIDE 25: Thank you ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2), '谢谢', font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(3.8), Inches(2), Pt(4), ORANGE)
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8), 'OpenClaw · The AI that actually does things',
    font_size=22, color=RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), 'V3b · 2026年3月 · Claude Opus 4',
    font_size=14, color=RGBColor(0x66, 0x77, 0x88), alignment=PP_ALIGN.CENTER)

# Save
output = '/root/.openclaw/workspace/reports/openclaw-product-insight-v3b.pptx'
os.makedirs(os.path.dirname(output), exist_ok=True)
prs.save(output)
print(f'Saved: {output}')
print(f'Slides: {len(prs.slides)}')
print(f'Size: {os.path.getsize(output) / 1024:.1f} KB')
